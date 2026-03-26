#!/usr/bin/env python3
"""
Create a 3-slide presentation: Introduction to Valmiki Ramayana
with ISKCON-themed design and relevant images from the internet.
"""

import os
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Configuration ────────────────────────────────────────────────────
OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_FILE = os.path.join(OUTPUT_DIR, "Valmiki_Ramayana_Introduction.pptx")

# Colour palette (saffron / gold / maroon — traditional ISKCON Ramayana theme)
SAFFRON       = RGBColor(0xFF, 0x99, 0x33)  # warm saffron
DARK_MAROON   = RGBColor(0x80, 0x00, 0x20)  # deep maroon
GOLD          = RGBColor(0xDA, 0xA5, 0x20)  # goldenrod
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
CREAM         = RGBColor(0xFF, 0xF8, 0xE7)
DARK_BG       = RGBColor(0x2B, 0x0A, 0x0A)  # very dark maroon-black
LIGHT_SAFFRON = RGBColor(0xFF, 0xE0, 0xB2)

# Image URLs (publicly available wallpapers / art)
IMAGES = {
    # Slide 1 - Title: Ram Sita with Hanuman
    "slide1": "https://wallpapers.com/images/hd/ram-sita-and-hanuman-3818i2viiuqzo9o1.jpg",
    # Slide 2 - About Valmiki: Ram Sita colorful painted art
    "slide2": "https://wallpapers.com/images/hd/ram-sita-colorful-painted-art-rv0k2mh8h0dj1baj.jpg",
    # Slide 3 - Seven Kandas: Ram Sita Laxman Hanuman
    "slide3": "https://wallpapers.com/images/hd/ram-sita-laxman-hanuman-n755amso7kbkpm9b.jpg",
}


def download_image(url: str) -> BytesIO | None:
    """Download an image and return as BytesIO, or None on failure."""
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
                          "AppleWebKit/537.36 (KHTML, like Gecko) "
                          "Chrome/120.0.0.0 Safari/537.36",
            "Referer": "https://wallpapers.com/",
        }
        resp = requests.get(url, headers=headers, timeout=30)
        resp.raise_for_status()
        buf = BytesIO(resp.content)
        buf.seek(0)
        return buf
    except Exception as e:
        print(f"  ⚠ Could not download {url}: {e}")
        return None


def add_background(slide, color):
    """Set solid-colour background for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, transparency=0):
    """Add a filled rectangle shape (useful for overlay bands)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # no border
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=WHITE, alignment=PP_ALIGN.LEFT,
                 font_name="Georgia", line_spacing=1.3):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(4)
    # line spacing
    from pptx.oxml.ns import qn
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    lnSpc = pPr.find(qn('a:lnSpc'))
    if lnSpc is None:
        from lxml import etree
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', str(int(line_spacing * 100000)))
    return txBox


def add_multiline_text(slide, left, top, width, height, lines,
                       font_size=16, color=WHITE, font_name="Georgia",
                       alignment=PP_ALIGN.LEFT, bold=False, line_spacing=1.2,
                       bullet=False):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        display_text = f"• {line}" if bullet else line
        p.text = display_text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(6)
    return txBox


# ─── Build Presentation ──────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

print("Downloading images …")
images = {}
for key, url in IMAGES.items():
    print(f"  → {key}: {url}")
    images[key] = download_image(url)

# =====================================================================
# SLIDE 1 — Title Slide
# =====================================================================
print("Creating Slide 1: Title …")
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_background(slide1, DARK_BG)

# Place image on the right side
if images["slide1"]:
    slide1.shapes.add_picture(images["slide1"],
                              Inches(6.5), Inches(0.5),
                              Inches(6.3), Inches(6.5))

# Semi-transparent overlay band on left
add_rect(slide1, Inches(0), Inches(0), Inches(7), SLIDE_H, DARK_BG)

# Decorative top accent bar
add_rect(slide1, Inches(0), Inches(0), SLIDE_W, Inches(0.12), SAFFRON)

# Decorative bottom accent bar
add_rect(slide1, Inches(0), SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), GOLD)

# Small OM / decorative text
add_text_box(slide1, Inches(0.8), Inches(1.2), Inches(5), Inches(0.6),
             "॥ श्री सीतारामचन्द्राय नमः ॥",
             font_size=16, bold=False, color=GOLD,
             alignment=PP_ALIGN.LEFT, font_name="Arial Unicode MS")

# Title
add_text_box(slide1, Inches(0.8), Inches(2.0), Inches(5.5), Inches(1.5),
             "Introduction to\nValmiki Ramayana",
             font_size=44, bold=True, color=WHITE,
             alignment=PP_ALIGN.LEFT, font_name="Georgia")

# Subtitle
add_text_box(slide1, Inches(0.8), Inches(3.8), Inches(5.5), Inches(1.0),
             "The Adi Kavya — The First Poem of Humanity",
             font_size=22, bold=False, color=LIGHT_SAFFRON,
             alignment=PP_ALIGN.LEFT, font_name="Georgia")

# Brief quote / description
add_text_box(slide1, Inches(0.8), Inches(5.0), Inches(5.5), Inches(1.5),
             "Composed by Maharishi Valmiki, the Ramayana is\n"
             "one of the two great Indian epics (Itihasa),\n"
             "narrating the divine pastimes of Lord Sri Rama.",
             font_size=15, bold=False, color=CREAM,
             alignment=PP_ALIGN.LEFT, font_name="Georgia", line_spacing=1.4)

# Attribution
add_text_box(slide1, Inches(0.8), Inches(6.5), Inches(5.5), Inches(0.5),
             "Images courtesy: ISKCON Desire Tree & Hare Krishna Wallpapers",
             font_size=10, bold=False, color=GOLD,
             alignment=PP_ALIGN.LEFT, font_name="Arial")


# =====================================================================
# SLIDE 2 — About Maharishi Valmiki & The Epic
# =====================================================================
print("Creating Slide 2: About Valmiki & The Epic …")
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide2, DARK_BG)

# Decorative top accent bar
add_rect(slide2, Inches(0), Inches(0), SLIDE_W, Inches(0.12), SAFFRON)
add_rect(slide2, Inches(0), SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), GOLD)

# Image on left
if images["slide2"]:
    slide2.shapes.add_picture(images["slide2"],
                              Inches(0.4), Inches(0.8),
                              Inches(5.0), Inches(6.0))

# Content panel on right
add_rect(slide2, Inches(5.8), Inches(0.4), Inches(7.2), Inches(6.7), RGBColor(0x3B, 0x12, 0x12))

# Section Title
add_text_box(slide2, Inches(6.2), Inches(0.7), Inches(6.5), Inches(0.8),
             "Maharishi Valmiki & The Epic",
             font_size=32, bold=True, color=SAFFRON,
             alignment=PP_ALIGN.LEFT, font_name="Georgia")

# Horizontal rule
add_rect(slide2, Inches(6.2), Inches(1.5), Inches(3), Inches(0.04), GOLD)

# About Valmiki
add_text_box(slide2, Inches(6.2), Inches(1.7), Inches(6.3), Inches(0.5),
             "About Valmiki",
             font_size=20, bold=True, color=GOLD,
             alignment=PP_ALIGN.LEFT, font_name="Georgia")

valmiki_points = [
    "Known as the 'Adi Kavi' (the First Poet)",
    "Composed the first shloka (verse) in Sanskrit poetry",
    "Inspired by witnessing a hunter's cruelty to a bird",
    "Sage Narada urged him to write Lord Rama's story",
]
add_multiline_text(slide2, Inches(6.2), Inches(2.2), Inches(6.3), Inches(2.0),
                   valmiki_points, font_size=15, color=CREAM,
                   font_name="Georgia", bullet=True, line_spacing=1.3)

# About the Epic
add_text_box(slide2, Inches(6.2), Inches(4.0), Inches(6.3), Inches(0.5),
             "The Ramayana at a Glance",
             font_size=20, bold=True, color=GOLD,
             alignment=PP_ALIGN.LEFT, font_name="Georgia")

epic_points = [
    "24,000 verses (shlokas) in 7 books (Kandas)",
    "Written in the Anushtubh metre",
    "Narrates events of Treta Yuga",
    "Central theme: Dharma, devotion & ideal conduct",
    "Rama is the Supreme Lord descended as the ideal king",
]
add_multiline_text(slide2, Inches(6.2), Inches(4.5), Inches(6.3), Inches(2.5),
                   epic_points, font_size=15, color=CREAM,
                   font_name="Georgia", bullet=True, line_spacing=1.3)


# =====================================================================
# SLIDE 3 — The Seven Kandas (Books)
# =====================================================================
print("Creating Slide 3: The Seven Kandas …")
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide3, DARK_BG)

# Decorative bars
add_rect(slide3, Inches(0), Inches(0), SLIDE_W, Inches(0.12), SAFFRON)
add_rect(slide3, Inches(0), SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), GOLD)

# Full-width faded image at top
if images["slide3"]:
    slide3.shapes.add_picture(images["slide3"],
                              Inches(0), Inches(0.12),
                              SLIDE_W, Inches(3.5))

# Overlay for readability over image
add_rect(slide3, Inches(0), Inches(0.12), SLIDE_W, Inches(3.5),
         RGBColor(0x2B, 0x0A, 0x0A))

# Title over overlay
add_text_box(slide3, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "The Seven Kandas (Books) of Valmiki Ramayana",
             font_size=32, bold=True, color=SAFFRON,
             alignment=PP_ALIGN.CENTER, font_name="Georgia")

add_rect(slide3, Inches(4.5), Inches(1.3), Inches(4.5), Inches(0.04), GOLD)

# Seven Kandas in two columns
kandas_left = [
    ("1. Bala Kanda", "Birth & early life of Rama; Rama breaks Shiva's bow"),
    ("2. Ayodhya Kanda", "Rama's exile for 14 years; Bharata's devotion"),
    ("3. Aranya Kanda", "Forest adventures; Sita's abduction by Ravana"),
    ("4. Kishkindha Kanda", "Alliance with Sugriva & Hanuman's emergence"),
]

kandas_right = [
    ("5. Sundara Kanda", "Hanuman leaps to Lanka; finds Sita; burns Lanka"),
    ("6. Yuddha Kanda", "Great war; Ravana's defeat; Rama reunites with Sita"),
    ("7. Uttara Kanda", "Rama's rule (Rama Rajya); the ideal kingdom"),
]

y_start = Inches(1.7)
left_x = Inches(0.6)
right_x = Inches(7.0)
card_w = Inches(5.8)
card_h = Inches(1.15)
gap = Inches(0.15)

for i, (title, desc) in enumerate(kandas_left):
    y = y_start + i * (card_h + gap)
    # Card background
    add_rect(slide3, left_x, y, card_w, card_h, RGBColor(0x3B, 0x12, 0x12))
    add_text_box(slide3, left_x + Inches(0.2), y + Inches(0.1), card_w - Inches(0.4), Inches(0.4),
                 title, font_size=16, bold=True, color=SAFFRON, font_name="Georgia")
    add_text_box(slide3, left_x + Inches(0.2), y + Inches(0.5), card_w - Inches(0.4), Inches(0.6),
                 desc, font_size=13, bold=False, color=CREAM, font_name="Georgia")

for i, (title, desc) in enumerate(kandas_right):
    y = y_start + i * (card_h + gap)
    add_rect(slide3, right_x, y, card_w, card_h, RGBColor(0x3B, 0x12, 0x12))
    add_text_box(slide3, right_x + Inches(0.2), y + Inches(0.1), card_w - Inches(0.4), Inches(0.4),
                 title, font_size=16, bold=True, color=SAFFRON, font_name="Georgia")
    add_text_box(slide3, right_x + Inches(0.2), y + Inches(0.5), card_w - Inches(0.4), Inches(0.6),
                 desc, font_size=13, bold=False, color=CREAM, font_name="Georgia")

# Bottom message
add_text_box(slide3, Inches(0.5), Inches(6.7), Inches(12), Inches(0.5),
             "\"Sitaayaah Charitam Mahat\" — Valmiki Ramayana is truly the noble story of Sita",
             font_size=14, bold=True, color=GOLD,
             alignment=PP_ALIGN.CENTER, font_name="Georgia")

# Attribution
add_text_box(slide3, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
             "Images courtesy: ISKCON Desire Tree & Hare Krishna Wallpapers  |  Jai Sri Ram!",
             font_size=10, bold=False, color=LIGHT_SAFFRON,
             alignment=PP_ALIGN.CENTER, font_name="Arial")


# ─── Save ─────────────────────────────────────────────────────────────
prs.save(OUTPUT_FILE)
print(f"\n✅ Presentation saved: {OUTPUT_FILE}")
